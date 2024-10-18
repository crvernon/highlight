import io
import os
import importlib

from docxtpl import DocxTemplate
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN

import streamlit as st

import highlight as hlt

if "config" not in st.session_state:
    st.session_state["config"] = hlt.read_config("config.toml")


if "reduce_document" not in st.session_state:
    st.session_state.reduce_document = False

if "content_dict" not in st.session_state:
    st.session_state.content_dict = {}

# parameters for word document
if "title_response" not in st.session_state:
    st.session_state.title_response = None

if "subtitle_response" not in st.session_state:
    st.session_state.subtitle_response = None

if "photo" not in st.session_state:
    st.session_state.photo = None

if "photo_link" not in st.session_state:
    st.session_state.photo_link = None

if "photo_site_name" not in st.session_state:
    st.session_state.photo_site_name = None

if "image_caption" not in st.session_state:
    st.session_state.image_caption = None

if "science_response" not in st.session_state:
    st.session_state.science_response = None

if "impact_response" not in st.session_state:
    st.session_state.impact_response = None

if "summary_response" not in st.session_state:
    st.session_state.summary_response = None

if "funding" not in st.session_state:
    st.session_state.funding = None

if "citation" not in st.session_state:
    st.session_state.citation = None

if "related_links" not in st.session_state:
    st.session_state.related_links = None

# additional word doc content that is not in the template
if "figure_response" not in st.session_state:
    st.session_state.figure_response = None

if "figure_caption" not in st.session_state:
    st.session_state.figure_caption = None

if "caption_response" not in st.session_state:
    st.session_state.caption_response = None

if "output_file" not in st.session_state:
    st.session_state.output_file = None

# parameters for the ppt slide
if "objective_response" not in st.session_state:
    st.session_state.objective_response = None

if "approach_response" not in st.session_state:
    st.session_state.approach_response = None

if "ppt_impact_response" not in st.session_state:
    st.session_state.ppt_impact_response = None

if "figure_recommendation" not in st.session_state:
    st.session_state.figure_recommendation = None

if "citation" not in st.session_state:
    st.session_state.citation = None

if "search_phrase" not in st.session_state:
    st.session_state.search_phrase = None

if "point_of_contact" not in st.session_state:
    st.session_state.point_of_contact = None

if "project_dict" not in st.session_state:
    st.session_state.project_dict = {
        "IM3": "Jennie Rice\nIM3 Principal Investigator\njennie.rice@pnnl.gov",
        "GCIMS": "Marshall Wise\nGCIMS Principal Investigator\nmarshall.wise@pnnl.gov",
        "COMPASS-GLM": "Robert Hetland\nCOMPASS-GLM Principal Investigator\nrobert.hetland@pnnl.gov",
        "ICoM": "Ian Kraucunas\nICoM Principal Investigator\nian.kraucunas@pnnl.gov",
        "Puget Sound": "Ning Sun\nPuget Sound Scoping and Pilot Study Principal Investigator\nning.sun@pnnl.gov",
        "Other": "First and Last Name\nCorresponding Project Name with POC Credentials\nEmail Address",
    }

# Force responsive layout for columns also on mobile
st.write(
    """<style>
    [data-testid="column"] {
        width: calc(50% - 1rem);
        flex: 1 1 calc(50% - 1rem);
        min-width: calc(50% - 1rem);
    }
    </style>""",
    unsafe_allow_html=True,
)

# Render streamlit page
st.title("Research Highlight Generator")

st.markdown(
    (
        "This app uses a Large Language Model (LLM) of your choosing to generate "
        + " formatted research highlight content from an input file."
    )
)

llm_id = st.selectbox(
    label="Select your model:", options=hlt.llm.list_llms(st.session_state["config"])
)
st.session_state.provider, st.session_state.model = llm_id.split(": ")
st.session_state.max_allowable_tokens = st.session_state.config["llm"][
    st.session_state.provider
][st.session_state.model]["max_allowable_tokens"]

st.session_state.client = hlt.llm.get_llm(
    st.session_state.provider,
    max_context_length=st.session_state.max_allowable_tokens,
)

st.markdown("### Upload file to process:")
uploaded_file = st.file_uploader(
    label="### Select PDF or text file to upload",
    type=["pdf", "txt"],
    help="Select PDF or text file to upload",
)

if uploaded_file is not None:

    if uploaded_file.type == "text/plain":
        content_dict = hlt.read_text(uploaded_file)

    elif uploaded_file.type == "application/pdf":
        content_dict = hlt.read_pdf(uploaded_file)

    st.session_state.output_file = uploaded_file.name

    st.code(
        f"""File specs:\n
    - Number of pages:  {content_dict['n_pages']}
    - Number of characters:  {content_dict['n_characters']}
    - Number of words: {content_dict['n_words']}
    - Number of tokens: {content_dict['n_tokens']}
    """
    )

    if content_dict["n_tokens"] > st.session_state.max_allowable_tokens:
        msg = f"""
    The number of tokens in your document exceeds the maximum allowable tokens.
    This will cause your queries to fail.
    The queries account for the number of tokens in a prompt + the number of tokens in your document.

    Maximum allowable token count: {st.session_state.max_allowable_tokens}

    Your documents token count: {content_dict['n_tokens']}

    Token deficit: {content_dict['n_tokens'] - st.session_state.max_allowable_tokens}
    """
        st.error(msg, icon="🚨")

        st.session_state.reduce_document = st.radio(
            """Would you like me to attempt to reduce the size of
        your document by keeping only relevant information?
        If so, I will give you a file to download with the content
        so you only have to do this once.
        If you choose to go through with this, it may take a while
        to process, usually on the order of 15 minutes for a 20K token
        document.
        Alternatively, you can copy and paste the contents that you
        know are of interest into a text file and upload that
        instead.

        """,
            ("Yes", "No"),
        )

    # word document content
    st.markdown("### Content to fill in Word document template:")

    # title section
    title_container = st.container()
    title_container.markdown("##### Generate title from text content")

    # title criteria
    title_container.markdown(
        """
    The title should meet the following criteria:
    - No colons are allowed in the output.
    - Should pique the interest of the reader while still being somewhat descriptive.
    - Be understandable to a general audience.
    - Should be only once sentence.
    - Should have a maximum length of 10 words.
    """
    )

    title_container.markdown("Set desired temperature:")

    # title slider
    title_temperature = title_container.slider(
        "Title Temperature", 0.0, 1.0, 0.2, label_visibility="collapsed"
    )

    # build container content
    if title_container.button("Generate Title"):

        st.session_state.title_response = hlt.generate_content(
            client=st.session_state.client,
            container=title_container,
            content=content_dict["content"],
            prompt_name="title",
            result_title="Title Result:",
            max_tokens=50,
            temperature=title_temperature,
            box_height=50,
            max_allowable_tokens=st.session_state.max_allowable_tokens,
            model=st.session_state.model,
        )

    else:
        if st.session_state.title_response is not None:
            title_container.markdown("Title Result:")
            title_container.text_area(
                label="Title Result:",
                value=st.session_state.title_response,
                label_visibility="collapsed",
                height=50,
            )

    # subtitle section
    subtitle_container = st.container()
    subtitle_container.markdown("##### Generate subtitle from text content")

    # subtitle criteria
    subtitle_container.markdown(
        """
    The subtitle should meet the following criteria:
    - Be an extension of and related to, but not directly quote, the title.
    - Provide information that will make the audience want to find out more about the research.
    - Do not use more than 155 characters including spaces.
    """
    )

    subtitle_container.markdown("Set desired temperature:")

    # subtitle slider
    subtitle_temperature = subtitle_container.slider(
        "Subtitle Temperature", 0.0, 1.0, 0.5, label_visibility="collapsed"
    )

    # build container content
    if subtitle_container.button("Generate Subtitle"):

        if st.session_state.title_response is None:
            st.write(
                "Please generate a Title first.  Subtitle generation considers the title response."
            )
        else:

            st.session_state.subtitle_response = hlt.generate_content(
                client=st.session_state.client,
                container=subtitle_container,
                content=content_dict["content"],
                prompt_name="subtitle",
                result_title="Subtitle Result:",
                max_tokens=100,
                temperature=subtitle_temperature,
                box_height=50,
                additional_content=st.session_state.title_response,
                max_word_count=100,
                min_word_count=75,
                max_allowable_tokens=st.session_state.max_allowable_tokens,
                model=st.session_state.model,
            )

    else:
        if st.session_state.subtitle_response is not None:
            subtitle_container.markdown("Subtitle Result:")
            subtitle_container.text_area(
                label="Subtitle Result:",
                value=st.session_state.subtitle_response,
                label_visibility="collapsed",
                height=50,
            )

    # science section
    science_container = st.container()
    science_container.markdown("##### Generate science summary from text content")

    # science criteria
    science_container.markdown(
        """
    **GOAL**:  Describe the scientific results for a non-expert, non-scientist audience.

    The description should meet the following criteria:
    - Answer what the big challenge in this field of science is that the research addresses.
    - State what the key finding is.
    - Explain the science, not the process.
    - Be understandable to a high school senior or college freshman.
    - Use short sentences and succinct words.
    - Avoid technical terms if possible.  If technical terms are necessary, define them.
    - Provide the necessary context so someone can have a very basic understanding of what you did.
    - Start with topics that the reader already may know and move on to more complex ideas.
    - Use present tense.
    - In general, the description should speak about the research or researchers in first person.
    - Use a minimum of 75 words and a maximum of 100 words.
    """
    )

    science_container.markdown("Set desired temperature:")

    # slider
    science_temperature = science_container.slider(
        "Science Summary Temperature", 0.0, 1.0, 0.3, label_visibility="collapsed"
    )

    # build container content
    if science_container.button("Generate Science Summary"):
        st.session_state.science_response = hlt.generate_content(
            client=st.session_state.client,
            container=science_container,
            content=content_dict["content"],
            prompt_name="science",
            result_title="Science Summary Result:",
            max_tokens=200,
            temperature=science_temperature,
            box_height=250,
            max_word_count=100,
            min_word_count=75,
            max_allowable_tokens=st.session_state.max_allowable_tokens,
            model=st.session_state.model,
        )

    else:
        if st.session_state.science_response is not None:
            science_container.markdown("Science Summary Result:")
            science_container.text_area(
                label="Science Summary Result:",
                value=st.session_state.science_response,
                label_visibility="collapsed",
                height=250,
            )

    # impact section
    impact_container = st.container()
    impact_container.markdown("##### Generate impact summary from text content")

    impact_container.markdown(
        """
    **GOAL**: Describe the impact of the research to a non-expert, non-scientist audience.

    The description should meet the following criteria:
    - Answer why the findings presented are important, i.e., what problem the research is trying to solve.
    - Answer if the finding is the first of its kind.
    - Answer what was innovative or distinct about the research.
    - Answer what the research enables other scientists in your field to do next.
    - Include other scientific fields potentially impacted.
    - Be understandable to a high school senior or college freshman.
    - Use short sentences and succinct words.
    - Avoid technical terms if possible.  If technical terms are necessary, define them.
    - Use present tense.
    - In general, the description should speak about the research or researchers in first person.
    - Use a minimum of 75 words and a maximum of 100 words.
    """
    )

    impact_container.markdown("Set desired temperature:")

    # slider
    impact_temperature = impact_container.slider(
        "Impact Summary Temperature", 0.0, 1.0, 0.0, label_visibility="collapsed"
    )

    # build container content
    if impact_container.button("Generate Impact Summary"):
        st.session_state.impact_response = hlt.generate_content(
            client=st.session_state.client,
            container=impact_container,
            content=content_dict["content"],
            prompt_name="impact",
            result_title="Impact Summary Result:",
            max_tokens=700,
            temperature=impact_temperature,
            box_height=250,
            max_word_count=100,
            min_word_count=75,
            max_allowable_tokens=st.session_state.max_allowable_tokens,
            model=st.session_state.model,
        )

    else:
        if st.session_state.impact_response is not None:
            impact_container.markdown("Impact Summary Result:")
            impact_container.text_area(
                label="Impact Summary Result:",
                value=st.session_state.impact_response,
                label_visibility="collapsed",
                height=250,
            )

    # general summary section
    summary_container = st.container()
    summary_container.markdown("##### Generate general summary from text content")

    summary_container.markdown(
        """
    **GOAL**: Generate a general summary of the current research.

    The summary should meet the following criteria:
    - Should relay key findings and value.
    - The summary should be still accessible to the non-specialist but may be more technical if necessary.
    - Do not mention the names of institutions.
    - If there is a United States Department of Energy Office of Science user facility involved, such as NERSC, you can mention the user facility.
    - Should be 1 or 2 paragraphs detailing the research.
    - Use present tense.
    - In general, the description should speak about the research or researchers in first person.
    - Use no more than 200 words.
    """
    )

    summary_container.markdown("Set desired temperature:")

    # slider
    summary_temperature = summary_container.slider(
        "General Summary Temperature", 0.0, 1.0, 0.3, label_visibility="collapsed"
    )

    # build container content
    if summary_container.button("Generate General Summary"):
        st.session_state.summary_response = hlt.generate_content(
            client=st.session_state.client,
            container=summary_container,
            content=content_dict["content"],
            prompt_name="summary",
            result_title="General Summary Result:",
            max_tokens=700,
            temperature=summary_temperature,
            box_height=400,
            max_word_count=200,
            min_word_count=100,
            max_allowable_tokens=st.session_state.max_allowable_tokens,
            model=st.session_state.model,
        )

    else:
        if st.session_state.summary_response is not None:
            summary_container.markdown("General Summary Result:")
            summary_container.text_area(
                label="General Summary Result:",
                value=st.session_state.summary_response,
                label_visibility="collapsed",
                height=400,
            )

    # figure recommendations section
    figure_container = st.container()
    figure_container.markdown(
        "##### Generate figure search string recommendations from the general summary"
    )
    figure_container.markdown("Set desired temperature:")

    # slider
    figure_temperature = figure_container.slider(
        "Figure Recommendations Temperature",
        0.0,
        1.0,
        0.9,
        label_visibility="collapsed",
    )

    # build container content
    if figure_container.button("Generate Figure Recommendations"):

        if st.session_state.summary_response is None:
            st.write("Please generate a general summary first.")
        else:
            st.session_state.figure_response = hlt.generate_content(
                client=st.session_state.client,
                container=figure_container,
                content=st.session_state.summary_response,
                prompt_name="figure",
                result_title="Figure Recommendations Result:",
                max_tokens=200,
                temperature=figure_temperature,
                box_height=200,
                max_allowable_tokens=st.session_state.max_allowable_tokens,
                model=st.session_state.model,
            )

    else:
        if st.session_state.figure_response is not None:

            figure_container.markdown("Figure Recommendations Result:")
            figure_container.text_area(
                label="Figure Recommendations Result:",
                value=st.session_state.figure_response,
                label_visibility="collapsed",
                height=200,
            )

    figure_summary_container = st.container()
    figure_summary_container.markdown(
        "##### Generate a figure caption that summarizes the work generally to use with the artistic photo above"
    )

    # slider
    figure_summary_container.markdown("Set desired temperature:")
    figure_summary_temperature = figure_summary_container.slider(
        "Figure Caption Temperature", 0.0, 1.0, 0.1, label_visibility="collapsed"
    )

    # build container content
    if figure_summary_container.button("Generate Figure Caption"):

        if st.session_state.summary_response is None:
            st.write("Please generate a general summary first.")
        else:
            st.session_state.figure_caption = hlt.generate_content(
                client=st.session_state.client,
                container=figure_summary_container,
                content=st.session_state.summary_response,
                prompt_name="figure_caption",
                result_title="Figure Caption Result:",
                max_tokens=300,
                temperature=figure_temperature,
                box_height=200,
                max_allowable_tokens=st.session_state.max_allowable_tokens,
                model=st.session_state.model,
            ).replace('"', "")

    else:
        if st.session_state.figure_caption is not None:
            figure_container.markdown("Figure Caption Result:")
            figure_container.text_area(
                label="Figure Caption Result:",
                value=st.session_state.figure_caption,
                label_visibility="collapsed",
                height=200,
            )

    # citation recommendations section
    citation_container = st.container()
    citation_container.markdown("##### Citation for the paper in Chicago style")

    if citation_container.button("Generate Citation"):
        st.session_state.citation = hlt.generate_content(
            client=st.session_state.client,
            container=citation_container,
            content=content_dict["content"],
            prompt_name="citation",
            result_title="",
            max_tokens=300,
            temperature=0.0,
            box_height=200,
            max_allowable_tokens=st.session_state.max_allowable_tokens,
            model=st.session_state.model,
        ).replace('"', "")

    else:
        if st.session_state.citation is not None:
            citation_container.text_area(
                label="Citation",
                value=st.session_state.citation,
                label_visibility="collapsed",
                height=200,
            )

    # funding recommendations section
    funding_container = st.container()
    funding_container.markdown("##### Funding statement from the paper")

    if funding_container.button("Generate funding statement"):
        st.session_state.funding = hlt.generate_content(
            client=st.session_state.client,
            container=funding_container,
            content=content_dict["content"],
            prompt_name="funding",
            result_title="",
            max_tokens=300,
            temperature=0.0,
            box_height=200,
            max_allowable_tokens=st.session_state.max_allowable_tokens,
            model=st.session_state.model,
        ).replace('"', "")

    else:
        if st.session_state.funding is not None:
            funding_container.text_area(
                label="Funding statement",
                value=st.session_state.funding,
                label_visibility="collapsed",
                height=200,
            )

    # point of contact box
    poc_container = st.container()
    poc_container.markdown("##### Point of contact for the research by project")

    # select the POC information from the dropdown
    st.session_state.point_of_contact = st.session_state.project_dict[
        poc_container.selectbox(
            label="Select the project who funded the work:",
            options=[
                "COMPASS-GLM",
                "GCIMS",
                "ICoM",
                "IM3",
                "Puget Sound",
                "Other",
            ],
        )
    ]

    poc_container.write("What will be written to the document as the point of contact:")
    poc_parts = st.session_state.point_of_contact.split("\n")
    poc_container.success(
        f"""
        {poc_parts[0]}\n
        {poc_parts[1]}\n
        {poc_parts[2]}\n
        """
    )

    export_container = st.container()
    export_container.markdown("##### Export Word document with new content when ready")

    # template parameters
    word_parameters = {
        "title": st.session_state.title_response,
        "subtitle": st.session_state.subtitle_response,
        "photo": st.session_state.photo,
        "photo_link": st.session_state.photo_link,
        "photo_site_name": st.session_state.photo_site_name,
        "image_caption": st.session_state.figure_caption,
        "science": st.session_state.science_response,
        "impact": st.session_state.impact_response,
        "summary": st.session_state.summary_response,
        "funding": st.session_state.funding,
        "citation": st.session_state.citation,
        "related_links": st.session_state.related_links,
        "point_of_contact": st.session_state.point_of_contact,
    }

    # template word document
    word_template_file = importlib.resources.files("highlight.data").joinpath(
        "highlight_template.docx"
    )
    template = DocxTemplate(word_template_file)

    template.render(word_parameters)
    bio = io.BytesIO()
    template.save(bio)
    if template:
        export_container.download_button(
            label="Export Word Document",
            data=bio.getvalue(),
            file_name="modified_template.docx",
            mime="docx",
        )

    # power point slide content
    st.markdown("### Content to fill in PowerPoint template:")

    # objective section
    objective_container = st.container()
    objective_container.markdown("##### Generate objective summary from text content")

    objective_container.markdown(
        """
    **GOAL**:  Generate one sentence stating the core purpose of the study.

    The sentence should meet the following criteria:
    - Use active verbs for the start of each point.
    - Use present tense.
    - Do not include methodology related to statistical, technological, and theory based
    """
    )

    objective_container.markdown("Set desired temperature:")

    # slider
    objective_temperature = objective_container.slider(
        "Objective Temperature", 0.0, 1.0, 0.3, label_visibility="collapsed"
    )

    # build container content
    if objective_container.button("Generate Objective"):
        st.session_state.objective_response = hlt.generate_content(
            client=st.session_state.client,
            container=objective_container,
            content=content_dict["content"],
            prompt_name="objective",
            result_title="Objective Result:",
            max_tokens=300,
            temperature=objective_temperature,
            box_height=250,
            max_allowable_tokens=st.session_state.max_allowable_tokens,
            model=st.session_state.model,
        )

    else:
        if st.session_state.objective_response is not None:
            objective_container.markdown("Objective Result:")
            objective_container.text_area(
                label="Objective Result:",
                value=st.session_state.objective_response,
                label_visibility="collapsed",
                height=250,
            )

    # approach section
    approach_container = st.container()
    approach_container.markdown("##### Generate approach summary from text content")

    approach_container.markdown(
        """
    **GOAL**:  Clearly and concisely state in 2-3 short points how this work accomplished the stated objective from a methodolgocial perspecive.
    - Based off of the objective summary
    - Only include methodology including but not limited to: statistical, technological, and theory based approaches.
    - Use a different action verb to start sentences than what is used to begin the objective statement.
    - Use active verbs for the start of each point.
    - Use present tense.
    """
    )

    approach_container.markdown("Set desired temperature:")

    # slider
    approach_temperature = approach_container.slider(
        "Approach Temperature", 0.0, 1.0, 0.1, label_visibility="collapsed"
    )

    # build container content
    if approach_container.button("Generate Approach"):
        st.session_state.approach_response = hlt.generate_content(
            client=st.session_state.client,
            container=approach_container,
            content=content_dict["content"],
            prompt_name="approach",
            result_title="Approach Result:",
            max_tokens=300,
            temperature=approach_temperature,
            box_height=250,
            additional_content=st.session_state.objective_response,
            max_allowable_tokens=st.session_state.max_allowable_tokens,
            model=st.session_state.model,
        )

    else:
        if st.session_state.approach_response is not None:
            approach_container.markdown("Approach Result:")
            approach_container.text_area(
                label="Approach Result:",
                value=st.session_state.approach_response,
                label_visibility="collapsed",
                height=250,
            )

    # power point impact section
    ppt_impact_container = st.container()
    ppt_impact_container.markdown("##### Generate impact points from text content")

    ppt_impact_container.markdown(
        """
    **GOAL**:  Clearly and concisely state in 3 points the key results and outcomes from this research.
    - State what the results indicate.
    - Include results that may be considered profound or surprising.
    - Each point should be 1 concise sentence.
    - Use present tense.
    """
    )

    ppt_impact_container.markdown("Set desired temperature:")

    # slider
    ppt_impact_temperature = ppt_impact_container.slider(
        "Impact Points Temperature", 0.0, 1.0, 0.1, label_visibility="collapsed"
    )

    # build container content
    if ppt_impact_container.button("Generate Impact Points"):
        st.session_state.ppt_impact_response = hlt.generate_content(
            client=st.session_state.client,
            container=ppt_impact_container,
            content=content_dict["content"],
            prompt_name="ppt_impact",
            result_title="Impact Points Result:",
            max_tokens=300,
            temperature=ppt_impact_temperature,
            box_height=250,
            max_allowable_tokens=st.session_state.max_allowable_tokens,
            model=st.session_state.model,
        )

    else:
        if st.session_state.ppt_impact_response is not None:
            ppt_impact_container.markdown("Impact Points Result:")
            ppt_impact_container.text_area(
                label="Impact Points Result:",
                value=st.session_state.ppt_impact_response,
                label_visibility="collapsed",
                height=250,
            )

    # power point figure selection section
    ppt_figure_selection = st.container()
    ppt_figure_selection.markdown("##### Select a representative figure from the paper")

    ppt_figure_selection.markdown(
        """
    **GOAL**:  What figure best represents the high impact content that can be easily understood by a non-technical, non-scientifc audience.

    Limit the response to:
    1. The figure name as it is written in the text,
    2. An explanation of why it was chosen,
    3. And what the figure is about in less than 50 words.
    """
    )

    ppt_figure_selection.markdown("Set desired temperature:")

    # slider
    ppt_figure_selection_temperature = ppt_figure_selection.slider(
        "Figure recommendation Temperature", 0.0, 1.0, 0.2, label_visibility="collapsed"
    )

    # build container content
    if ppt_figure_selection.button("Generate Figure Recommendation"):
        st.session_state.figure_recommendation = hlt.generate_content(
            client=st.session_state.client,
            container=ppt_figure_selection,
            content=content_dict["content"],
            prompt_name="figure_choice",
            result_title="Figure Recommendation Result:",
            max_tokens=300,
            temperature=ppt_figure_selection_temperature,
            box_height=250,
            max_allowable_tokens=st.session_state.max_allowable_tokens,
            model=st.session_state.model,
        )

    else:
        if st.session_state.figure_recommendation is not None:
            ppt_figure_selection.markdown("Figure Recommendation Result:")
            ppt_figure_selection.text_area(
                label="Figure Recommendation Result:",
                value=st.session_state.figure_recommendation,
                label_visibility="collapsed",
                height=250,
            )

    # Add PowerPoint export container at the end
    export_ppt_container = st.container()
    export_ppt_container.markdown(
        "##### Export PowerPoint Presentation with New Content"
    )

    if (
        "title_response" in st.session_state
        and "objective_response" in st.session_state
        and "ppt_impact_response" in st.session_state
        and "approach_response" in st.session_state
    ):

        if export_ppt_container.button("Export PowerPoint"):

            try:
                # Load the PowerPoint template
                ppt_template_file = importlib.resources.files(
                    "highlight.data"
                ).joinpath("highlight_template.pptx")
                prs = Presentation(ppt_template_file)

                # Split the impact and approach responses into bullet points (assuming they are separated by newlines)
                impact_points = st.session_state.ppt_impact_response.split("\n")
                approach_points = st.session_state.approach_response.split("\n")

                # Iterate over all slides to find the text boxes labeled "impact_0", "impact_1", "impact_2"
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            # Handle title insertion and maintain font size and bold
                            if "title" in shape.text_frame.text:
                                shape.text_frame.text = st.session_state.title_response

                                # Ensure font size and bold settings are maintained for each paragraph
                                for paragraph in shape.text_frame.paragraphs:
                                    for run in paragraph.runs:
                                        run.font.size = Pt(
                                            24
                                        )  # Example size, adjust as needed
                                        run.font.bold = True  # Maintain bold
                                        run.alignment = PP_ALIGN.LEFT  # Align title

                            # Handle citation insertion and maintain font size and bold
                            if "citation" in shape.text_frame.text:
                                shape.text_frame.text = st.session_state.citation

                                # Ensure font size and bold settings are maintained for each paragraph
                                for paragraph in shape.text_frame.paragraphs:
                                    for run in paragraph.runs:
                                        run.font.size = Pt(
                                            11
                                        )  # Example size for citation, adjust as needed
                                        run.font.bold = (
                                            False  # Citation typically isn't bold
                                        )
                                        run.alignment = PP_ALIGN.LEFT  # Align citation

                            if shape.text_frame.text == "objective_0":
                                # Set the text of the text box to the objective response
                                shape.text_frame.text = (
                                    st.session_state.objective_response
                                )

                                # Optional: Adjust font size and alignment for the objective
                                for paragraph in shape.text_frame.paragraphs:
                                    paragraph.font.size = Pt(13)  # Set font size
                                    paragraph.alignment = PP_ALIGN.LEFT  # Set alignment

                            # Handle approach bullet points
                            if "approach_0" in shape.text_frame.text:
                                # Clear existing paragraphs in the text frame
                                shape.text_frame.clear()

                                # Add bullet points for approach
                                for i, approach_point in enumerate(
                                    approach_points[:3]
                                ):  # Only take the first 3 approach points
                                    p = shape.text_frame.add_paragraph()
                                    p.text = approach_point
                                    p.level = 0  # This sets it as a bullet point
                                    p.font.size = Pt(
                                        13
                                    )  # Adjust bullet point font size
                                    p.alignment = PP_ALIGN.LEFT  # Align bullet points

                            # Handle the impact bullet points in the same text box
                            if "impact_0" in shape.text_frame.text:
                                # Clear the existing paragraphs
                                shape.text_frame.clear()

                                # Add bullet points for impact
                                for i, impact_point in enumerate(
                                    impact_points[:3]
                                ):  # Only take the first 3 impact points
                                    p = shape.text_frame.add_paragraph()
                                    p.text = impact_point
                                    p.level = 0  # This sets it as a bullet point
                                    p.font.size = Pt(
                                        13
                                    )  # Adjust bullet point font size
                                    p.alignment = PP_ALIGN.LEFT  # Align bullet points

                # Save the modified presentation to a BytesIO object
                ppt_io = io.BytesIO()
                prs.save(ppt_io)
                ppt_io.seek(0)

                # Provide a download button for the user
                export_ppt_container.download_button(
                    label="Export PowerPoint Presentation",
                    data=ppt_io,
                    file_name="modified_highlight_template.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                )

                export_ppt_container.success(
                    "PowerPoint presentation generated successfully!", icon="✅"
                )

            except Exception as e:
                export_ppt_container.error(
                    f"An error occurred while generating the PowerPoint: {e}", icon="🚨"
                )

    else:
        export_ppt_container.error(
            "Please generate the objective and impact responses before exporting.",
            icon="⚠️",
        )
